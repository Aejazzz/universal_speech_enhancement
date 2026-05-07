#!/usr/bin/env python3
"""
Build a minimal Google-inspired .pptx for Universal Speech Enhancement.

Usage:
    python scripts/build_presentation.py [--out docs/Presentation_Universal_Speech_Enhancement.pptx]

Requires: python-pptx
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

# Google Material-ish palette (minimal)
C_TITLE = "202124"  # near-black
C_BODY = "5F6368"  # gray
C_ACCENT = "1A73E8"  # Google blue
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx import Presentation


def _rgb(hex6: str) -> RGBColor:
    return RGBColor(int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16))


def _set_run_font(run, size_pt: int, bold: bool = False, color_hex: str = C_BODY):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = "Segoe UI"
    run.font.color.rgb = _rgb(color_hex)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # subtle top bar
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.12))  # rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(C_ACCENT)
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(11.5), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    _set_run_font(run, 36, bold=True, color_hex=C_TITLE)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(12)
    run2 = p2.add_run()
    run2.text = subtitle
    _set_run_font(run2, 16, bold=False, color_hex=C_BODY)

    foot = slide.shapes.add_textbox(Inches(0.7), Inches(6.8), Inches(11.0), Inches(0.5))
    ft = foot.text_frame
    fp = ft.paragraphs[0]
    fr = fp.add_run()
    fr.text = "github.com/Aejazzz/universal_speech_enhancement"
    _set_run_font(fr, 11, color_hex=C_ACCENT)


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bar = slide.shapes.add_shape(1, 0, Inches(0.35), prs.slide_width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(C_ACCENT)
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.65), Inches(0.55), Inches(11.5), Inches(0.85))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_run_font(r, 28, bold=True, color_hex=C_TITLE)

    body = slide.shapes.add_textbox(Inches(0.85), Inches(1.45), Inches(11.2), Inches(5.5))
    btf = body.text_frame
    btf.word_wrap = True
    for i, line in enumerate(bullets):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.space_after = Pt(8)
        para.level = 0
        run = para.add_run()
        run.text = line
        _set_run_font(run, 15, color_hex=C_BODY)


def add_two_column_slide(prs: Presentation, title: str, left_title: str, left_lines: list[str], right_title: str, right_lines: list[str]) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bar = slide.shapes.add_shape(1, 0, Inches(0.35), prs.slide_width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(C_ACCENT)
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.65), Inches(0.55), Inches(11.5), Inches(0.85))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_run_font(r, 28, bold=True, color_hex=C_TITLE)

    def col(x: float, t: str, lines: list[str]):
        tt = slide.shapes.add_textbox(Inches(x), Inches(1.35), Inches(5.4), Inches(0.4))
        ttf = tt.text_frame
        tp = ttf.paragraphs[0]
        tr = tp.add_run()
        tr.text = t
        _set_run_font(tr, 14, bold=True, color_hex=C_ACCENT)
        box = slide.shapes.add_textbox(Inches(x), Inches(1.75), Inches(5.4), Inches(5.0))
        btf = box.text_frame
        for i, line in enumerate(lines):
            para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            para.space_after = Pt(6)
            run = para.add_run()
            run.text = line
            _set_run_font(run, 13, color_hex=C_BODY)

    col(0.65, left_title, left_lines)
    col(6.85, right_title, right_lines)


def add_image_slide(prs: Presentation, title: str, img_path: Path, max_h_inches: float = 4.8) -> bool:
    if not img_path.is_file():
        return False
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bar = slide.shapes.add_shape(1, 0, Inches(0.35), prs.slide_width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(C_ACCENT)
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.65), Inches(0.55), Inches(11.5), Inches(0.75))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_run_font(r, 24, bold=True, color_hex=C_TITLE)

    # scale image to fit width ~11 in, max height
    try:
        pic = slide.shapes.add_picture(str(img_path), Inches(1.0), Inches(1.35))
    except Exception:
        return False
    # aspect ratio preserve
    h, w = pic.height, pic.width
    max_w = Inches(11.0)
    max_h = Inches(max_h_inches)
    scale = min(max_w / w, max_h / h)
    pic.width = int(w * scale)
    pic.height = int(h * scale)
    pic.left = int((prs.slide_width - pic.width) / 2)
    pic.top = Inches(1.25)
    return True


def add_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]]) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bar = slide.shapes.add_shape(1, 0, Inches(0.35), prs.slide_width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(C_ACCENT)
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.65), Inches(0.55), Inches(11.5), Inches(0.75))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_run_font(r, 24, bold=True, color_hex=C_TITLE)

    nrows, ncols = len(rows) + 1, len(headers)
    tbl = slide.shapes.add_table(nrows, ncols, Inches(0.65), Inches(1.35), Inches(12.0), Inches(0.45 * nrows + 0.2)).table

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                _set_run_font(run, 11, bold=True, color_hex=C_TITLE)

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    _set_run_font(run, 11, color_hex=C_BODY)


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument(
        "--out",
        default=str(ROOT / "docs" / "Presentation_Universal_Speech_Enhancement.pptx"),
        help="Output .pptx path",
    )
    ap.add_argument(
        "--plots-dir",
        default=str(ROOT / "outputs" / "reports" / "demo_blind_eval" / "plots"),
        help="Directory with PNG plots (optional)",
    )
    args = ap.parse_args()
    out_path = Path(args.out)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    plots_dir = Path(args.plots_dir)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Universal Speech Enhancement",
        "Policy learning · multi-expert routing · DNSMOS-driven selection\n"
        "Frozen WavLM + dynamic speculate-and-measure · FastAPI + React",
    )

    add_bullet_slide(
        prs,
        "Problem & approach",
        [
            "Real-world noisy speech: one fixed denoiser rarely fits all distortions.",
            "Mixture-of-experts: neural models (DeepFilterNet3, optional MossFormer2 / Resemble) plus classical FFT baselines.",
            "Trained router (WavLM + MLP) suggests expert and strength; dynamic router scores every active candidate with DNSMOS P.835 (+ UTMOS when reliable).",
            '"Do no harm" margin keeps BYPASS when improvement is marginal.',
        ],
    )

    add_bullet_slide(
        prs,
        "End-to-end pipeline",
        [
            "Load & resample → 16 kHz mono",
            "Preprocess: DC removal, 60 Hz HPF, ITU-R BS.1770 loudness (−23 LUFS)",
            "Distortion analyzer (CRNN on log-mel) → 6-D features for policy",
            "Policy: frozen WavLM-base-plus, multi-layer pool → expert, strength, refine",
            "Dynamic routing (optional): sweep strengths per expert → rank by DNSMOS blend",
            "Brick-wall limiter −1 dBFS → metrics (DNSMOS, PESQ, STOI, SI-SDR) + plots",
        ],
    )

    add_two_column_slide(
        prs,
        "Architecture snapshot",
        "Stack",
        [
            "Frontend: React, Tailwind, Vite",
            "Backend: FastAPI, uvicorn",
            "Core: PyTorch, speechmos (DNSMOS), transformers (WavLM)",
        ],
        "Artifacts per run",
        [
            "enhanced.wav, metrics.json, routing.json",
            "Waveform / spectrogram / policy probability PNGs",
            "Full dynamic candidate leaderboard in API",
        ],
    )

    add_table_slide(
        prs,
        "Blind evaluation — 100 clips (URGENT validation noisy set)",
        ["Metric", "Value"],
        [
            ["Files processed", "100"],
            ["Mean ΔDNSMOS (OVRL)", "+0.72"],
            ["Mean ΔSIG", "+0.56"],
            ["Mean ΔBAK", "+1.54"],
            ["Chosen = top rank", "100%"],
            ["Pareto-dominated picks", "0"],
        ],
    )

    add_table_slide(
        prs,
        "Routing distribution (mean ΔDNSMOS by chosen expert)",
        ["Expert", "Clips", "Mean ΔOVRL", "Mean ΔSIG", "Mean ΔBAK"],
        [
            ["DeepFilterNet3", "57", "+0.74", "+0.47", "+1.60"],
            ["NoiseReduce", "26", "+0.58", "+0.27", "+1.40"],
            ["WienerFilter", "11", "+1.01", "+1.51", "+1.64"],
            ["SpectralSubtraction", "4", "+0.55", "+0.84", "+1.38"],
            ["SpectralGate", "2", "+0.89", "+1.14", "+1.49"],
        ],
    )

    # Figures from eval (if present)
    for title, fname in [
        ("Results: routing distribution", "01_routing_distribution.png"),
        ("Results: expert leaderboard (mean deltas)", "03_expert_leaderboard.png"),
        ("Results: noisy vs enhanced DNSMOS (OVRL)", "04_orig_vs_enhanced_scatter.png"),
    ]:
        p = plots_dir / fname
        if not add_image_slide(prs, title, p):
            add_bullet_slide(
                prs,
                title + " (run plot script)",
                [
                    f"Plot not found: {p}",
                    "Generate: python scripts/plot_blind_results.py --eval-dir outputs/reports/demo_blind_eval",
                ],
            )

    add_bullet_slide(
        prs,
        "Live demo",
        [
            "Windows: pwsh -File scripts/start_demo.ps1",
            "Backend http://127.0.0.1:8001/health · Frontend http://127.0.0.1:5173",
            "Upload .wav / .flac / .mp3 → Enhance Audio → scorecard, leaderboard, A/B audio",
            "Batch eval: python scripts/eval_blind_batch.py --input-dir <noisy_folder> --resume",
        ],
    )

    add_bullet_slide(
        prs,
        "Notes for Q&A",
        [
            "UTMOS22 may fall back to heuristic if torch.hub conflicts with speechmos; router then uses DNSMOS-only blend (utmos_reliable=False).",
            "SOTA comparison bars in plots are illustrative — cite papers/leaderboards for publication.",
            "See README.md and DEMO.md for full architecture, API schema, and training pipeline.",
        ],
    )

    prs.save(str(out_path))
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    main()
